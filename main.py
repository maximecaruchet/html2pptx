#!/usr/bin/env python3

import http.server
import io
import os
import re
import requests
import urllib.parse
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT

SHORT_TEXT_LIMIT_CHARS = 75
TITLE_FONT_PT = 75
SLIDE_BLANK_LAYOUT = 6
SLIDE_WIDTH_INCHES = 10
SLIDE_HEIGHT_INCHES = 7.5
SLIDE_SMALL_MARGIN_INCHES = 0.25
COLUMN_MARGIN_INCHES = 0.1
HEIGHT_MARGIN_INCHES = 0.1

# Init configuration default values
debug_logs = False
debug_slides = False
server_port = 8080

# Parse configuration environment variables
html2pptx_debug_logs = os.getenv('HTML2PPTX_DEBUG_LOGS', "false")
html2pptx_debug_slides = os.getenv('HTML2PPTX_DEBUG_SLIDES', "false")
html2pptx_server_port = os.getenv('HTML2PPTX_PORT', "8080")

# Set configuration variables depending on environment variables
if html2pptx_debug_logs.lower() == "true":
    debug_logs = True
if html2pptx_debug_slides.lower() == "true":
    debug_slides = True
server_port = int(html2pptx_server_port)

# Print configuration
print("Configuration:")
print("debug_logs:", debug_logs)
print("debug_slides:", debug_slides)
print("server_port:", server_port)


def html_to_pptx(url, css_selector):
    r = requests.get(url)
    url_string = r.text
    slides = html_to_slides(url_string, css_selector)
    prs_bytes_stream = slides_to_pptx(slides)
    return prs_bytes_stream


def html_to_slides(html_string, css_selector):
    soup = BeautifulSoup(html_string, 'html.parser')
    useful_content = soup.select(css_selector)
    slides = []
    for parent_content_tag in useful_content[0].children:
        if parent_content_tag.name is not None:
            slide_content = html_to_slide(parent_content_tag)
            slides.append(slide_content)
    return slides


def html_to_slide(parent_tag):
    return parse_tag_contents(parent_tag)


def parse_tag_contents(tag):
    tag_data = []
    for children_content_tag in tag.children:
        # Go through all children tags
        if children_content_tag.name is not None:
            # Just handle valid tags
            if children_content_tag.name == "img":
                # If we have an image, get the "src" link
                tag_data.append("img_src:" + children_content_tag["src"])
            elif children_content_tag.string is not None:
                # If we have only one string, return it
                if children_content_tag.string.strip() != "":
                    tag_data.append(children_content_tag.string.strip())
            else:
                # Get direct text elements from tag even if there are children elements with text inside
                # (but do not get the text from the children)
                direct_tag_strings = children_content_tag.find_all(string=True, recursive=False)
                sanitized_direct_tag_strings = []
                for string in direct_tag_strings:
                    sanitized_string = string.strip()
                    if re.match("^\\[if mso \\| IE\\].*", sanitized_string):
                        sanitized_string = ""
                    if sanitized_string != "":
                        sanitized_direct_tag_strings.append(sanitized_string)

                # Get direct text elements from tag even if there are children elements with text inside
                # (this time, we get the text from the children)
                recursive_tag_strings = children_content_tag.find_all(string=True, recursive=True)
                sanitized_recursive_tag_strings = []
                for string in recursive_tag_strings:
                    sanitized_string = string.strip()
                    if re.match("^\\[if mso \\| IE\\].*", sanitized_string):
                        sanitized_string = ""
                    if sanitized_string != "":
                        sanitized_recursive_tag_strings.append(sanitized_string)

                # If we have some direct text elements, then we are in a case of some formatted text nested within other
                # text tags, then just extract the whole text (direct and from children), return it,
                # and stop the recursion by going directly to the next element
                if len(sanitized_direct_tag_strings) > 0:
                    tag_data.append(" ".join(sanitized_recursive_tag_strings))
                    continue

                # If we are not in the case of nested text, just do a recursive call
                # to get the contents of the children tag
                tag_data.extend(parse_tag_contents(children_content_tag))
    return tag_data


def slides_to_pptx(slides):
    prs = Presentation()
    for slide in slides:
        if debug_logs:
            print("============================================ NEW SLIDE ============================================")
        fill_slide(prs, slide)
        if debug_logs:
            print("============================================ END SLIDE ============================================")
    prs_bytes_stream = io.BytesIO()
    prs.save(prs_bytes_stream)

    # Test to save bytes stream directly to file
    # Uncomment if needed
    # with open("test.pptx", 'wb') as out:
    #     out.write(prs_bytes_stream.getvalue())
    return prs_bytes_stream


def fill_slide(prs, slide):
    # Init default count values
    image_count = 0
    max_chars_in_strings = 0

    # Determine number of images and max string length
    for slide_data in slide:
        img_found = re.search("^img_src:(.*)$", slide_data)
        if img_found:
            image_count += 1
        else:
            if len(slide_data) > max_chars_in_strings:
                max_chars_in_strings = len(slide_data)

    # Determine if the slide is empty (no images and, max string length = 0, wich means no images and no text)
    empty_slide = image_count == 0 and max_chars_in_strings == 0
    if empty_slide:
        return

    # Add a slide when not empty
    prs_slide_layout = prs.slide_layouts[SLIDE_BLANK_LAYOUT]
    prs_slide = prs.slides.add_slide(prs_slide_layout)

    # Determine if we are in a column layout or not
    # We have a column layout if we have more than 1 image, at least 1 text block
    # and if the longest text is < SHORT_TEXT_LIMIT_CHARS chars
    with_multiple_images = image_count > 1
    with_short_texts = max_chars_in_strings != 0 and max_chars_in_strings <= SHORT_TEXT_LIMIT_CHARS
    column_layout = with_multiple_images and with_short_texts

    # Init base data
    images_array = []
    text_array = []
    column_text_array = []

    # Parse slide and separate images from text
    for slide_data in slide:
        image_found = re.search("^img_src:(.*)$", slide_data)
        if column_layout:
            # Special handling of text if we are in a column layout
            if image_found:
                # Be sure to always associate an image with some text below
                if len(images_array) == 0 and len(column_text_array) > 0:
                    # If text has been found without an image at the beginning, add an "empty" placeholder image
                    images_array.append("empty")
                    text_array.append(column_text_array)
                    column_text_array = []
                elif len(images_array) > 0:
                    # When some images have already been found, just add associated text
                    # to the previous image (empty or not)
                    text_array.append(column_text_array)
                    column_text_array = []
                images_array.append(image_found.group(1))
            else:
                # Add text to the current column
                column_text_array.append(slide_data)
        else:
            # Default handling when not in column layout
            if image_found:
                images_array.append(image_found.group(1))
            else:
                column_text_array.append(slide_data)

    if column_layout:
        # Always add the final text (empty or not) in column layout
        text_array.append(column_text_array)
    else:
        # Add text only if some text has been found (there will only be one big column in this case)
        if len(column_text_array) > 0:
            text_array.append(column_text_array)

    # Determine available space in slide and column layout
    available_slide_width_inches = SLIDE_WIDTH_INCHES - 2*SLIDE_SMALL_MARGIN_INCHES \
        - (len(images_array)-1)*COLUMN_MARGIN_INCHES
    column_width_inches = available_slide_width_inches
    if len(images_array) > 0:
        column_width_inches = available_slide_width_inches/len(images_array)

    # Init default values to compute image heights
    max_image_height_inches = 0
    images_heights_inches = []

    for index, image_link in enumerate(images_array):
        if debug_logs:
            print("IMAGE LINK:", image_link)

        if image_link == "empty":
            # If we have an "empty" placeholder image, ignore it but add a height of "0" to the image heights array
            images_heights_inches.append(0)
            continue

        # Determine image position
        top = Inches(SLIDE_SMALL_MARGIN_INCHES)
        left = Inches(SLIDE_SMALL_MARGIN_INCHES + index*column_width_inches
                      + index*COLUMN_MARGIN_INCHES)

        # Download image and add it to the slide
        image_req = requests.get(image_link)
        image_bytes = io.BytesIO(image_req.content)
        image_box = prs_slide.shapes.add_picture(image_bytes, left, top)

        # Determine image ratio and resize image (with aspect ratio preserved) to fit it in the slide if it is too large
        ratio = image_box.width.inches / image_box.height.inches
        if image_box.width.inches > SLIDE_WIDTH_INCHES - 2*SLIDE_SMALL_MARGIN_INCHES:
            image_box.width = Inches(SLIDE_WIDTH_INCHES - 2*SLIDE_SMALL_MARGIN_INCHES)
            image_box.height = Inches(image_box.width.inches / ratio)
        if image_box.height.inches > SLIDE_HEIGHT_INCHES - 2*SLIDE_SMALL_MARGIN_INCHES:
            image_box.height = Inches(SLIDE_HEIGHT_INCHES - 2*SLIDE_SMALL_MARGIN_INCHES)
            image_box.width = Inches(image_box.height.inches * ratio)

        # Center image horizontally if only one image is found
        if len(images_array) == 1:
            horizontal_image_center_inches = image_box.width.inches / 2
            slide_horizontal_center_inches = SLIDE_WIDTH_INCHES / 2
            left_horizontal_centered_inches = slide_horizontal_center_inches - horizontal_image_center_inches
            if left_horizontal_centered_inches < SLIDE_SMALL_MARGIN_INCHES:
                left_horizontal_centered_inches = SLIDE_SMALL_MARGIN_INCHES
            image_box.left = Inches(left_horizontal_centered_inches)

            # Center image vertically if this one image is alone with no text column
            if len(text_array) == 0:
                vertical_image_center_inches = image_box.height.inches / 2
                slide_vertical_center_inches = SLIDE_HEIGHT_INCHES / 2
                top_vertical_centered_inches = slide_vertical_center_inches - vertical_image_center_inches
                if top_vertical_centered_inches < SLIDE_SMALL_MARGIN_INCHES:
                    top_vertical_centered_inches = SLIDE_SMALL_MARGIN_INCHES
                image_box.top = Inches(top_vertical_centered_inches)

        images_heights_inches.append(image_box.height.inches)
        if image_box.height.inches > max_image_height_inches:
            max_image_height_inches = image_box.height.inches

    for index, text_column in enumerate(text_array):
        # For every text column, init default position and size
        left = Inches(SLIDE_SMALL_MARGIN_INCHES)
        width = Inches(SLIDE_WIDTH_INCHES - 2*SLIDE_SMALL_MARGIN_INCHES)
        top = Inches(SLIDE_SMALL_MARGIN_INCHES)
        height = Inches(SLIDE_HEIGHT_INCHES - 2*SLIDE_SMALL_MARGIN_INCHES)

        if len(images_array) > 0:
            # Override some default values if we have images
            top = Inches(SLIDE_SMALL_MARGIN_INCHES + max_image_height_inches + HEIGHT_MARGIN_INCHES)
            height = Inches(SLIDE_HEIGHT_INCHES - 2*SLIDE_SMALL_MARGIN_INCHES
                            - max_image_height_inches - HEIGHT_MARGIN_INCHES)

        if column_layout:
            # Column layout gets the final override if enabled
            left = Inches(SLIDE_SMALL_MARGIN_INCHES + index*column_width_inches
                          + index*COLUMN_MARGIN_INCHES)
            width = Inches(column_width_inches)
            top = Inches(SLIDE_SMALL_MARGIN_INCHES + images_heights_inches[index] + HEIGHT_MARGIN_INCHES)
            height = Inches(SLIDE_HEIGHT_INCHES - 2*SLIDE_SMALL_MARGIN_INCHES
                            - images_heights_inches[index] - HEIGHT_MARGIN_INCHES)

        # Create the text box
        text_box = prs_slide.shapes.add_textbox(left, top, width, height)

        if debug_slides:
            # Fill the shape with red in debug mode
            fill = text_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 0, 0)

        # Create the text frame inside the text box and configure it
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        if column_layout:
            # Change vertical anchor in column layout mode
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Add every string in the column to the text frame
        for text in text_column:
            # Determine if the text is a title (no images + only this text alone)
            is_title = len(images_array) == 0 and len(text_array) == 1 and len(text_column) == 1

            if debug_logs:
                print(text)

            # Fill the text frame
            paragraph = text_frame.paragraphs[0]
            if paragraph.text == "":
                if is_title:
                    # title format
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    paragraph.font.size = Pt(TITLE_FONT_PT)
                paragraph.text = text
            else:
                paragraph = text_frame.add_paragraph()
                if is_title:
                    # title format
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    paragraph.font.size = Pt(TITLE_FONT_PT)
                paragraph.text = text


class Html2pptx(http.server.BaseHTTPRequestHandler):
    def do_GET(self):
        # Handle GET requests

        # Set headers
        self.send_response(200)
        self.send_header("Content-type", "text/html")
        self.end_headers()

        # Send index.html page
        with open("index.html", 'rb') as out:
            self.wfile.write(out.read())

    def do_POST(self):
        # Handle POST requests

        # Retrieve and decode POST query data
        content_length = int(self.headers['Content-Length'])
        post_data = self.rfile.read(content_length)
        decoded_post_data = urllib.parse.parse_qs(post_data.decode("utf-8"))
        if debug_logs:
            print("decoded_post_data[\"url\"][0]", decoded_post_data["url"][0])
            print("decoded_post_data[\"selector\"][0]", decoded_post_data["selector"][0])

        # Translate HTML to PPTX, retrieves presentation bytes stream
        prs_bytes_stream = html_to_pptx(decoded_post_data["url"][0], decoded_post_data["selector"][0])

        # Set headers to download the PPTX file
        self.send_response(200)
        self.send_header("Content-Type", 'application/octet-stream')
        self.send_header("Content-Disposition", 'attachment; filename="presentation.pptx"')
        # This is some unused example code which uses content-length header when transferring a file
        # Since it seems to work here, we won't use it, but the code below will stay
        # here in case we need to use and modify it
        # Source:
        # https://stackoverflow.com/questions/18543640/how-would-i-create-a-python-web-server-that-downloads-a-file-on-any-get-request
        # fs = os.fstat(f.fileno())
        # self.send_header("Content-Length", str(fs.st_size))
        self.end_headers()

        # Send the PPTX presentation
        # Use getvalue() instead of read() with BytesIO to avoid problems
        # Source:
        # https://stackoverflow.com/questions/46981529/why-does-saving-a-presentation-to-a-file-like-object-produce-a-blank-presentatio?noredirect=1&lq=1
        self.wfile.write(prs_bytes_stream.getvalue())


# Setup and start HTTP server with custom Html2pptx handler
server_address = ("", server_port)
httpd = http.server.HTTPServer(server_address, Html2pptx)
print("Serving at port:", server_port)
httpd.serve_forever()
